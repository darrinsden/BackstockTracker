"""
Generates BackstockTracker-IT-intro.pptx at the repo root.

Audience: head of IT. Goal: explain what the app is, how it's built,
what it depends on, and what IT needs to support / approve to keep it
on solid ground (Apple Dev seats, iCloud container ownership, Firebase
project for the cross-platform migration, Drive permissions).
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# ---- Palette ---------------------------------------------------------------

NAVY = RGBColor(0x0F, 0x2A, 0x44)
TEAL = RGBColor(0x1F, 0x6F, 0x8B)
ACCENT = RGBColor(0xE8, 0x8B, 0x1A)
LIGHT = RGBColor(0xF5, 0xF6, 0xF8)
GREY = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


def add_band(slide, top, height, color):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SLIDE_W, height)
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = color
    return band


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=NAVY, align=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, bullets, *,
                size=18, color=NAVY):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "•  " + b
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_header(slide, title, subtitle=None):
    add_band(slide, 0, Inches(1.1), NAVY)
    add_text(
        slide, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.6),
        title, size=30, bold=True, color=WHITE,
    )
    if subtitle:
        add_text(
            slide, Inches(0.5), Inches(0.66), Inches(12.3), Inches(0.4),
            subtitle, size=14, color=RGBColor(0xCF, 0xDC, 0xE7),
        )
    # Thin accent bar under header.
    add_band(slide, Inches(1.1), Inches(0.06), ACCENT)


def add_footer(slide, page, total):
    add_text(
        slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
        "Backstock Tracker  ·  IT introduction  ·  May 2026",
        size=10, color=GREY,
    )
    add_text(
        slide, Inches(11.8), Inches(7.05), Inches(1.3), Inches(0.3),
        f"{page} / {total}", size=10, color=GREY,
    )


# ---- Slide builders --------------------------------------------------------

SLIDES = []


def slide(fn):
    SLIDES.append(fn)
    return fn


@slide
def title_slide(s):
    add_band(s, 0, SLIDE_H, NAVY)
    add_band(s, Inches(3.4), Inches(0.08), ACCENT)
    add_text(
        s, Inches(0.7), Inches(2.5), Inches(12), Inches(1.2),
        "Backstock Tracker", size=60, bold=True, color=WHITE,
    )
    add_text(
        s, Inches(0.7), Inches(3.55), Inches(12), Inches(0.6),
        "iOS field tool for Jacent Area Managers",
        size=24, color=RGBColor(0xCF, 0xDC, 0xE7),
    )
    add_text(
        s, Inches(0.7), Inches(4.25), Inches(12), Inches(0.5),
        "Tracks scanned product credits against the $149.99 per-box limit.",
        size=16, color=RGBColor(0xB8, 0xC8, 0xD6),
    )
    add_text(
        s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.4),
        "Prepared for the Head of IT  ·  Darrin Horn  ·  May 2026",
        size=12, color=RGBColor(0x9F, 0xB2, 0xC2),
    )


@slide
def problem(s):
    add_header(s, "The problem we're solving",
               "What AMs do today, and why it's manual")
    add_bullets(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(5.4), [
        "Area Managers scan returned / clearance product into boxes during store visits.",
        "Every box has a hard $149.99 credit limit; over-limit boxes need TM email approval.",
        "Before this app: paper logs + a calculator + a phone call to the TM.",
        "Math errors meant rejected boxes, time on rework, and lost credits.",
        "Per-store catalog scoping (same UPC, different prices at Target vs. Walmart) made paper-tracking error-prone.",
    ], size=18)


@slide
def what_it_does(s):
    add_header(s, "60-second product tour",
               "What the AM does, screen by screen")
    add_bullets(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(5.4), [
        "Sign in once — picks identity from the roster CSV (no passwords).",
        "Scan UPC with the camera (VisionKit) — instant catalog lookup, running total.",
        "Audible confirm / error tones; red banner for not-in-catalog and wrong-store.",
        "At $149.99 the Submit button becomes Request Approval — opens a pre-filled email to the TM.",
        "Submitted box auto-uploads to the team feed; history is filterable per-store and per-area.",
        "Pick-list flow lets any AM in the area mark items to be pulled from another AM's box.",
    ], size=17)


@slide
def users_scale(s):
    add_header(s, "Users & operating scale",
               "Who's on it, who supports it")
    add_bullets(s, Inches(0.6), Inches(1.5), Inches(6), Inches(5.4), [
        "Jacent AMs nationwide (single role today).",
        "Roster sourced from one Drive CSV — ~hundreds of active AMs.",
        "Each AM works one Area; Areas roll up into Territories; Territories own a TM (approver).",
        "Catalog ~10s of thousands of SKUs, refreshed daily.",
        "Currently on TestFlight; build 17 is the latest.",
    ], size=16)
    add_bullets(s, Inches(6.8), Inches(1.5), Inches(6.2), Inches(5.4), [
        "iPhone-only (camera scanning), iOS 17+.",
        "Offline-capable scanning; sync-when-online for team feed.",
        "No back-office UI — Drive CSVs are the admin surface.",
        "Single maintainer today (Darrin) — this presentation is in part to widen that.",
    ], size=16)


@slide
def architecture(s):
    add_header(s, "Architecture at a glance",
               "SwiftUI front, SwiftData on-device, CloudKit for team sync")
    # Three column boxes.
    cols = [
        ("Client (iOS)", [
            "SwiftUI + SwiftData",
            "Single-file app (~10K LOC)",
            "VisionKit DataScanner",
            "AVFoundation for audio cues",
            "Zero third-party SPM deps",
        ]),
        ("Reference data (Drive)", [
            "area_managers.csv  ·  roster",
            "catalog.csv  ·  UPC → price/name",
            "stores.csv  ·  area scoping",
            "territory_managers.csv  ·  TM email",
            "Atomic full-replace on each sync",
        ]),
        ("Team sync (CloudKit)", [
            "iCloud public DB",
            "Record type: BackstockSession",
            "Anonymous records — no PII",
            "Area-scoped queries",
            "Optimistic + retry on launch",
        ]),
    ]
    x = Inches(0.5)
    for title, items in cols:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5),
                                 Inches(4.1), Inches(5.2))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = TEAL
        box.line.width = Pt(1.5)
        add_text(s, x, Inches(1.7), Inches(4.1), Inches(0.5), title,
                 size=18, bold=True, color=TEAL,
                 align=None)
        add_bullets(s, x + Inches(0.15), Inches(2.3),
                    Inches(3.85), Inches(4.2), items, size=14)
        x += Inches(4.2)


@slide
def data_sources(s):
    add_header(s, "Reference data: 4 Drive CSVs",
               "Each file is the canonical source for one domain")
    rows = [
        ("area_managers.csv", "employeeNumber, firstName, lastName, territory, area, email"),
        ("catalog.csv", "upc, name, price, commodity, store, retailPrice, rank"),
        ("stores.csv", "store, storeNumber, area, shortName"),
        ("territory_managers.csv", "territory, email"),
    ]
    y = Inches(1.6)
    add_text(s, Inches(0.6), y, Inches(4), Inches(0.4), "File",
             size=14, bold=True, color=GREY)
    add_text(s, Inches(4.8), y, Inches(8), Inches(0.4), "Schema",
             size=14, bold=True, color=GREY)
    y = Inches(2.05)
    for name, schema in rows:
        add_band(s, y - Inches(0.05), Inches(0.7), LIGHT)
        add_text(s, Inches(0.6), y, Inches(4), Inches(0.6), name,
                 size=16, bold=True, color=NAVY)
        add_text(s, Inches(4.8), y, Inches(8), Inches(0.6), schema,
                 size=14, color=NAVY)
        y += Inches(0.85)
    add_text(
        s, Inches(0.6), Inches(5.85), Inches(12.2), Inches(1.2),
        "All files are Google Drive “Anyone with link can view.” App accepts any "
        "Drive URL form; SyncService.normalizeSourceURL rewrites to the direct-"
        "download form before fetching. Each sync is atomic full-replace, so a "
        "bad row never half-applies.",
        size=13, color=GREY,
    )


@slide
def security(s):
    add_header(s, "Security & data residency",
               "What leaves the device, what doesn't")
    add_bullets(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(5.4), [
        "Reference CSVs: read-only fetch from Drive, anonymous (no Google auth).",
        "Submitted boxes: pushed to CloudKit public DB as anonymous records — no submitter identity, no employee number, no email.",
        "Local-only PII: the signed-in AM's name / email lives in SwiftData on-device, never in the cloud record.",
        "TM approval: sent through the user's own Mail.app — Jacent Exchange / GWS handles transport; we never relay through a third-party server.",
        "Camera: scoped to barcode scanning only (VisionKit DataScanner). Info.plist NSCameraUsageDescription pending before App Store release.",
        "No analytics SDKs, no crash-reporter SDKs, no third-party network endpoints other than Apple + Google Drive.",
    ], size=16)


@slide
def schema_migration(s):
    add_header(s, "On-device schema evolution",
               "How we ship model changes without nuking installs")
    add_bullets(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(5.4), [
        "SwiftData ModelContainer init runs through a BackstockMigrationPlan.",
        "Today: one version (BackstockSchemaV1) holding 8 @Model types.",
        "Future change → copy to V2, add a MigrationStage (.lightweight for additive, .custom for renames/splits).",
        "If migration ever fails for real disk/sandbox reasons, app routes to a StorageErrorView with a Copy-diagnostics affordance — no crash on first launch.",
        "Test-flight build numbers are bumped manually; archive → upload to App Store Connect.",
    ], size=17)


@slide
def cloudkit_state(s):
    add_header(s, "CloudKit team-sync — current state",
               "What's deployed, what's still required for production")
    add_bullets(s, Inches(0.6), Inches(1.5), Inches(6), Inches(5.4), [
        "Container: iCloud.com.jacent.BackstockTracker",
        "Record type: BackstockSession",
        "Idempotent on session UUID (no duplicate records)",
        "Optimistic cache + retryPending sweep on launch",
        "Read-side scrub for a known catalog-text corruption issue (chain-name leak)",
    ], size=15)
    add_bullets(s, Inches(6.8), Inches(1.5), Inches(6.2), Inches(5.4), [
        "TO DO before production:",
        "Make submittedAt Sortable, area Queryable",
        "Open BackstockSession Write + Delete to Authenticated",
        "Deploy schema from Development to Production",
        "Enable iCloud capability in Xcode signing",
    ], size=15)


@slide
def cross_platform(s):
    add_header(s, "Cross-platform readiness",
               "Why a Firestore migration is staged behind a flag")
    add_bullets(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(5.4), [
        "CloudKit ties team-sync to Apple devices. Android port is on the roadmap.",
        "FirestoreSyncService.swift is in the repo as scaffolding (not yet linked into the target).",
        "Public API mirrors CloudSyncService 1:1 — call sites can be flipped behind a useFirestore feature flag without touching the views.",
        "Plan: dual-write for one TestFlight cycle, then flip reads, then retire CloudKit writes.",
        "data-contract.md (repo root) is the canonical Firestore schema — same shape across iOS and the eventual Android client.",
    ], size=16)


@slide
def dependencies(s):
    add_header(s, "Dependencies & supply chain",
               "Deliberately small")
    add_bullets(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(5.4), [
        "Today: zero third-party SPM packages, zero CocoaPods. Apple frameworks only.",
        "Frameworks in use: SwiftUI, SwiftData, AVFoundation, VisionKit, Vision, MessageUI, CloudKit, BackgroundTasks.",
        "External services: Google Drive (CSV hosting, read-only), Apple CloudKit (team sync), Apple iCloud Auth (anonymous).",
        "Planned additions (Firestore migration): FirebaseFirestore + FirebaseAuth via SPM — first third-party deps. Documented exemption in CLAUDE.md.",
        "Build toolchain: Xcode 15+, ships to TestFlight from a single developer Mac today.",
    ], size=16)


@slide
def risks(s):
    add_header(s, "Known risks & open items",
               "What I want IT eyes on")
    add_bullets(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(5.4), [
        "Bus factor: single maintainer (me). Need a second engineer or contractor with Swift familiarity.",
        "Apple Dev account: currently personal; should be a Jacent organisation account for continuity.",
        "iCloud container is under the personal Dev account — same continuity concern. Migration to a Jacent container = re-key + roster reset.",
        "Drive CSV permissions: anyone-with-link. Acceptable for read-only catalog data; should be reviewed against Jacent data-handling policy.",
        "No formal observability: errors surface as in-app banners + manual TestFlight reports. Crashlytics-equivalent (or Apple's Xcode Organizer crashes) needs to be the bar before larger rollout.",
        "Pending Info.plist additions (NSCameraUsageDescription, BG task identifiers) before App Store submission.",
    ], size=15)


@slide
def roadmap(s):
    add_header(s, "Roadmap (next 2 quarters)",
               "Sequenced for risk reduction first, features second")
    add_bullets(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(5.4), [
        "Q2 2026 — CloudKit production deploy, App Store submission, NSCameraUsageDescription, organisation Dev account.",
        "Q2 2026 — Firestore scaffolding wired in behind a flag; dual-write one TestFlight cycle.",
        "Q3 2026 — Flip reads to Firestore; retire CloudKit writes; HistoryView filter chips.",
        "Q3 2026 — Android port (same Drive sources, same Firestore contract).",
        "Q3 2026 — Background-refresh sync for catalog/roster; in-app diagnostics export.",
    ], size=17)


@slide
def asks(s):
    add_header(s, "What I need from IT",
               "Concrete, sequenced asks")
    add_bullets(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(5.4), [
        "1.  Jacent-owned Apple Developer Program enrolment + transfer of the app + iCloud container.",
        "2.  Decision on Drive hosting — keep as-is (read-only anyone-with-link) or move to a service account / shared drive.",
        "3.  Approval to add Firebase (Firestore + Auth) as the cross-platform backend; an org-owned GCP project to host it.",
        "4.  Sign-off on data classification: anonymous record contents (store / store# / box / item names / prices / submittedAt) — confirm none is PII per Jacent policy.",
        "5.  A second engineer (Swift or generalist) for bus-factor and Android port.",
        "6.  Distribution choice for production: public App Store, or managed in-house via Apple Business Manager.",
    ], size=15)


@slide
def closing(s):
    add_band(s, 0, SLIDE_H, NAVY)
    add_band(s, Inches(3.4), Inches(0.08), ACCENT)
    add_text(
        s, Inches(0.7), Inches(2.6), Inches(12), Inches(1.2),
        "Questions?", size=54, bold=True, color=WHITE,
    )
    add_text(
        s, Inches(0.7), Inches(3.8), Inches(12), Inches(0.6),
        "Darrin Horn  ·  darrinsden@mac.com",
        size=20, color=RGBColor(0xCF, 0xDC, 0xE7),
    )
    add_text(
        s, Inches(0.7), Inches(4.3), Inches(12), Inches(0.6),
        "Repo: ~/jacent/BackstockTracker  ·  Branch: main",
        size=16, color=RGBColor(0xB8, 0xC8, 0xD6),
    )


# ---- Render ----------------------------------------------------------------

total = len(SLIDES)
for i, fn in enumerate(SLIDES, start=1):
    s = prs.slides.add_slide(BLANK)
    fn(s)
    if i > 1 and i < total:
        add_footer(s, i, total)

out = Path(__file__).resolve().parent.parent / "BackstockTracker-IT-intro.pptx"
prs.save(out)
print(f"Wrote {out}")
